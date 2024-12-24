from pptx import Presentation
from pptx.util import Cm
import datetime

path = "/home/daniel/PycharmProjects/Fiscalia/InformeViolaciones/"
template = "Template.pptx"
file = f'{datetime.date.today().strftime("%Y%m%d")} ActualizacionViolacion.pptx'

root = Presentation(f'{path}{template}')

slide2 = root.slides.add_slide(root.slide_layouts[0])
pic1 = slide2.shapes.add_picture(f'{path}ProdCD.png', Cm(1.2), Cm(0.1))

slide3 = root.slides.add_slide(root.slide_layouts[0])
pic2 = slide3.shapes.add_picture(f'{path}ProdSD.png', Cm(1.2), Cm(0.1))

slide4 = root.slides.add_slide(root.slide_layouts[0])
pic3 = slide4.shapes.add_picture(f'{path}ProdOASol.png', Cm(1.2), Cm(0.1))

slide5 = root.slides.add_slide(root.slide_layouts[0])
pic4 = slide5.shapes.add_picture(f'{path}ProdOACum.png', Cm(1.2), Cm(0.1))

root.save(f'{path}{file}')
